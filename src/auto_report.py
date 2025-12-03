"""
Core pipeline for automated report generation.
Ingests CSV, computes KPIs, creates charts, and generates PPTX presentation.
"""

import argparse
import logging
import os
from datetime import datetime
from pathlib import Path
from typing import Optional

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

try:
    from src.utils import (
        validate_csv_data,
        compute_kpis,
        create_time_series_chart,
        create_top_campaigns_chart,
        get_top_campaigns_list,
        llm_generate_summary,
        llm_generate_summary_with_groq
    )
except ImportError:
    # Fallback for direct script execution
    import sys
    from pathlib import Path
    sys.path.insert(0, str(Path(__file__).parent.parent))
    from src.utils import (
        validate_csv_data,
        compute_kpis,
        create_time_series_chart,
        create_top_campaigns_chart,
        get_top_campaigns_list,
        llm_generate_summary,
        llm_generate_summary_with_groq
    )

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


def load_and_filter_data(
    csv_path: Optional[str] = None,
    date_start: Optional[str] = None,
    date_end: Optional[str] = None,
    campaigns: Optional[list] = None,
    locations: Optional[list] = None,
    channels: Optional[list] = None,
    dataframe: Optional[pd.DataFrame] = None
) -> pd.DataFrame:
    """
    Load CSV data and apply filters.
    
    Args:
        csv_path: Path to CSV file
        date_start: Start date filter (YYYY-MM-DD)
        date_end: End date filter (YYYY-MM-DD)
        campaigns: List of campaign names to filter
        locations: List of locations to filter (if column exists)
        channels: List of channels to filter (if column exists)
        
    Returns:
        Filtered DataFrame
    """
    if dataframe is not None:
        df = dataframe.copy()
        logger.info(f"Filtering in-memory dataframe with {len(df)} rows")
    else:
        if not csv_path:
            raise ValueError("csv_path is required when dataframe is not provided.")
        logger.info(f"Loading data from {csv_path}")
        df = pd.read_csv(csv_path)
    
    # Validate data
    is_valid, error_msg = validate_csv_data(df)
    if not is_valid:
        raise ValueError(f"Data validation failed: {error_msg}")
    
    # Convert Date to datetime
    df['Date'] = pd.to_datetime(df['Date'])
    
    # Apply date filters
    if date_start:
        df = df[df['Date'] >= pd.to_datetime(date_start)]
    if date_end:
        df = df[df['Date'] <= pd.to_datetime(date_end)]
    
    # Apply campaign filter
    if campaigns:
        df = df[df['Campaign'].isin(campaigns)]
    
    # Apply location filter (if column exists)
    if locations and 'Location' in df.columns:
        df = df[df['Location'].isin(locations)]
    
    # Apply channel filter (if column exists)
    if channels and 'Channel' in df.columns:
        df = df[df['Channel'].isin(channels)]
    
    logger.info(f"Loaded {len(df)} rows after filtering")
    return df


def build_llm_context(df: pd.DataFrame, kpis: dict, top_campaigns: list, 
                       date_start: Optional[str], date_end: Optional[str]) -> dict:
    """
    Build context dictionary for LLM from aggregated data.
    
    Args:
        df: Filtered DataFrame
        kpis: Computed KPI dictionary
        top_campaigns: List of top campaign dictionaries
        date_start: Start date string
        date_end: End date string
        
    Returns:
        Context dictionary for LLM
    """
    # Build time window string
    if date_start and date_end:
        time_window = f"{date_start} to {date_end}"
    elif 'Date' in df.columns:
        df['Date'] = pd.to_datetime(df['Date'])
        min_date = df['Date'].min().strftime('%Y-%m-%d')
        max_date = df['Date'].max().strftime('%Y-%m-%d')
        time_window = f"{min_date} to {max_date}"
    else:
        time_window = "Unknown period"
    
    # Build totals
    totals = {
        'impressions': float(kpis.get('total_impressions', 0)),
        'clicks': float(kpis.get('total_clicks', 0)),
        'spend': float(kpis.get('total_spend', 0)),
        'visits': float(kpis.get('total_visits', 0))
    }
    
    # Build metrics
    metrics = {
        'avg_ctr': float(kpis.get('avg_ctr', 0)),
        'avg_cpc': float(kpis.get('avg_cpc', 0))
    }
    
    # Calculate ROI if possible (simplified: visits / spend)
    if totals['spend'] > 0:
        metrics['avg_roi'] = totals['visits'] / totals['spend']
    else:
        metrics['avg_roi'] = 0.0
    
    # Build top campaigns (limit to top 3)
    top_campaigns_list = []
    for campaign in top_campaigns[:3]:
        top_campaigns_list.append({
            'name': str(campaign.get('Campaign', '')),
            'visits': float(campaign.get('Visits', 0)),
            'ctr': float(campaign.get('CTR', 0)),
            'spend': float(campaign.get('Spend', 0))
        })
    
    context = {
        'time_window': time_window,
        'totals': totals,
        'metrics': metrics,
        'top_campaigns': top_campaigns_list
    }
    
    return context


def generate_report(
    output_path: str,
    csv_path: Optional[str] = None,
    date_start: Optional[str] = None,
    date_end: Optional[str] = None,
    campaigns: Optional[list] = None,
    locations: Optional[list] = None,
    channels: Optional[list] = None,
    use_llm: bool = False,
    use_groq: bool = False,
    use_cache: bool = True,
    temp_dir: str = "outputs",
    dataframe: Optional[pd.DataFrame] = None
) -> str:
    """
    Main pipeline: ingest data, compute KPIs, create charts, generate PPTX.
    
    Args:
        output_path: Path to save output PPTX
        csv_path: Path to input CSV file (optional if dataframe provided)
        date_start: Start date filter
        date_end: End date filter
        campaigns: Campaign filter list
        locations: Location filter list
        channels: Channel filter list
        use_llm: Whether to use LLM for summary generation
        use_cache: Cache toggle for Groq summary
        temp_dir: Temporary directory for chart images
        dataframe: Optional in-memory dataframe (bypasses csv_path)
        
    Returns:
        Path to generated PPTX file
    """
    logger.info("=" * 60)
    logger.info("Starting Automated Report Generation Pipeline")
    logger.info("=" * 60)
    
    # Step 1: Ingest and filter data
    logger.info("Step 1: Ingesting and filtering data")
    df = load_and_filter_data(
        csv_path=csv_path,
        date_start=date_start,
        date_end=date_end,
        campaigns=campaigns,
        locations=locations,
        channels=channels,
        dataframe=dataframe
    )
    
    if df.empty:
        raise ValueError("No data remaining after applying filters")
    
    # Step 2: Compute KPIs
    logger.info("Step 2: Computing KPIs")
    kpis = compute_kpis(df)
    
    # Step 3: Get top campaigns
    logger.info("Step 3: Analyzing top campaigns")
    top_campaigns = get_top_campaigns_list(df, top_n=5)
    
    # Step 4: Create visualizations
    logger.info("Step 4: Creating visualizations")
    os.makedirs(temp_dir, exist_ok=True)
    
    chart1_path = os.path.join(temp_dir, "time_series_chart.png")
    chart2_path = os.path.join(temp_dir, "top_campaigns_chart.png")
    
    create_time_series_chart(df, chart1_path)
    create_top_campaigns_chart(df, chart2_path, top_n=5)
    
    # Step 5: Generate executive summary
    logger.info("Step 5: Generating executive summary")
    
    if use_groq:
        # Build context for Groq
        context = build_llm_context(df, kpis, top_campaigns, date_start, date_end)
        summary, latency = llm_generate_summary_with_groq(context, use_cache=use_cache)
        if latency:
            logger.info(f"Groq summary generated in {latency:.2f}s")
    elif use_llm:
        # Use legacy LLM function (OpenAI/Groq fallback)
        summary = llm_generate_summary(kpis, top_campaigns, use_llm=True)
    else:
        # Use template
        from src.utils import _generate_template_summary
        summary = _generate_template_summary(kpis, top_campaigns)
    
    # Step 6: Create PPTX presentation
    logger.info("Step 6: Creating PPTX presentation")
    prs = create_presentation(kpis, top_campaigns, summary, chart1_path, chart2_path)
    
    # Save presentation
    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
    prs.save(output_path)
    
    logger.info(f"Report generated successfully: {output_path}")
    logger.info("=" * 60)
    
    return output_path


def create_presentation(
    kpis: dict,
    top_campaigns: list,
    summary: str,
    chart1_path: str,
    chart2_path: str
) -> Presentation:
    """
    Create a PowerPoint presentation with slides.
    
    Args:
        kpis: Dictionary of KPI metrics
        top_campaigns: List of top campaign dictionaries
        summary: Executive summary text
        chart1_path: Path to time series chart
        chart2_path: Path to top campaigns chart
        
    Returns:
        Presentation object
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Marketing Campaign Performance Report"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    timestamp_box = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.5))
    timestamp_frame = timestamp_box.text_frame
    timestamp_frame.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    timestamp_para = timestamp_frame.paragraphs[0]
    timestamp_para.font.size = Pt(18)
    timestamp_para.alignment = PP_ALIGN.CENTER
    
    # KPI snapshot on title slide
    kpi_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.6), Inches(7), Inches(2))
    kpi_frame = kpi_box.text_frame
    kpi_frame.word_wrap = True
    kpi_paragraphs = [
        f"Total Impressions: {kpis['total_impressions']:,.0f}",
        f"Total Clicks: {kpis['total_clicks']:,.0f}",
        f"Total Spend: ${kpis['total_spend']:,.2f}",
        f"Total Visits: {kpis['total_visits']:,.0f}",
        f"Avg CTR: {kpis['avg_ctr']:.2f}%",
        f"Avg CPC: ${kpis['avg_cpc']:.2f}"
    ]
    kpi_frame.text = "Key Metrics Snapshot"
    first_para = kpi_frame.paragraphs[0]
    first_para.font.size = Pt(20)
    first_para.font.bold = True
    first_para.space_after = Pt(10)
    
    for metric in kpi_paragraphs:
        p = kpi_frame.add_paragraph()
        p.text = f"• {metric}"
        p.font.size = Pt(16)
        p.space_after = Pt(4)
    
    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    summary_title = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
    summary_title_frame = summary_title.text_frame
    summary_title_frame.text = "Executive Summary"
    summary_title_frame.paragraphs[0].font.size = Pt(32)
    summary_title_frame.paragraphs[0].font.bold = True
    
    summary_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    summary_frame = summary_box.text_frame
    summary_frame.text = summary
    summary_frame.word_wrap = True
    for para in summary_frame.paragraphs:
        para.font.size = Pt(14)
        para.space_after = Pt(6)
    
    # Slide 3: Key Visuals (Two charts side-by-side)
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    visuals_title = slide3.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.4))
    visuals_title_frame = visuals_title.text_frame
    visuals_title_frame.text = "Key Visuals"
    visuals_title_frame.paragraphs[0].font.size = Pt(28)
    visuals_title_frame.paragraphs[0].font.bold = True
    
    # Chart 1 (left)
    slide3.shapes.add_picture(chart1_path, Inches(0.5), Inches(0.8), width=Inches(4.5))
    
    # Chart 2 (right)
    slide3.shapes.add_picture(chart2_path, Inches(5.2), Inches(0.8), width=Inches(4.5))
    
    # Slide 4: Top Campaigns
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    campaigns_title = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.4))
    campaigns_title_frame = campaigns_title.text_frame
    campaigns_title_frame.text = "Top 5 Campaigns"
    campaigns_title_frame.paragraphs[0].font.size = Pt(28)
    campaigns_title_frame.paragraphs[0].font.bold = True
    
    campaigns_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(6))
    campaigns_frame = campaigns_box.text_frame
    campaigns_frame.word_wrap = True
    
    for i, campaign in enumerate(top_campaigns, 1):
        text = f"{i}. {campaign['Campaign']}\n"
        text += f"   • Visits: {campaign['Visits']:,.0f}\n"
        text += f"   • Clicks: {campaign['Clicks']:,.0f}\n"
        text += f"   • Spend: ${campaign['Spend']:,.2f}\n"
        text += f"   • CTR: {campaign['CTR']:.2f}%\n\n"
        
        if i == 1:
            campaigns_frame.text = text
        else:
            p = campaigns_frame.add_paragraph()
            p.text = text
        
        campaigns_frame.paragraphs[-1].font.size = Pt(14)
        campaigns_frame.paragraphs[-1].space_after = Pt(8)
    
    # Slide 5: Appendix (Basic Stats)
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    appendix_title = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.4))
    appendix_title_frame = appendix_title.text_frame
    appendix_title_frame.text = "Appendix: Key Metrics"
    appendix_title_frame.paragraphs[0].font.size = Pt(28)
    appendix_title_frame.paragraphs[0].font.bold = True
    
    appendix_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(6))
    appendix_frame = appendix_box.text_frame
    appendix_frame.word_wrap = True
    
    stats_text = f"Total Impressions: {kpis['total_impressions']:,.0f}\n"
    stats_text += f"Total Clicks: {kpis['total_clicks']:,.0f}\n"
    stats_text += f"Total Spend: ${kpis['total_spend']:,.2f}\n"
    stats_text += f"Total Visits: {kpis['total_visits']:,.0f}\n"
    stats_text += f"Average CTR: {kpis['avg_ctr']:.2f}%\n"
    stats_text += f"Average CPC: ${kpis['avg_cpc']:.2f}\n"
    
    appendix_frame.text = stats_text
    for para in appendix_frame.paragraphs:
        para.font.size = Pt(16)
        para.space_after = Pt(10)
    
    return prs


def main():
    """Command-line interface for batch report generation."""
    parser = argparse.ArgumentParser(description='Generate automated marketing campaign report')
    parser.add_argument('--input', help='Path to input CSV file')
    parser.add_argument('--out', default='outputs/automated_report.pptx', help='Output PPTX path')
    parser.add_argument('--date-start', help='Start date filter (YYYY-MM-DD)')
    parser.add_argument('--date-end', help='End date filter (YYYY-MM-DD)')
    parser.add_argument('--campaigns', nargs='+', help='Campaign names to filter')
    parser.add_argument('--no-llm', action='store_true', help='Disable LLM summary generation')
    parser.add_argument('--use-groq', action='store_true', help='Use Groq AI for summary generation')
    parser.add_argument('--no-cache', action='store_true', help='Disable LLM response caching')
    parser.add_argument('--clear-llm-cache', action='store_true', help='Clear LLM cache and exit')
    parser.add_argument('--data-source', choices=['csv', 'sql', 'table'], default='csv',
                        help='Select data source type (default: csv)')
    parser.add_argument('--db-conn', help='SQLAlchemy connection string for SQL/table sources')
    parser.add_argument('--db-query', help='SQL query to execute when data-source=sql')
    parser.add_argument('--db-table', help='Table name to read when data-source=table')
    parser.add_argument('--db-limit', type=int, default=10000,
                        help='Row limit when loading a table (data-source=table)')
    
    args = parser.parse_args()
    
    # Handle cache clearing
    if args.clear_llm_cache:
        from src.utils import clear_llm_cache
        clear_llm_cache()
        print("✓ LLM cache cleared")
        return
    
    use_llm = not args.no_llm
    use_groq = args.use_groq
    use_cache = not args.no_cache
    
    df_source = None
    if args.data_source == 'csv':
        if not args.input:
            parser.error("--input is required when data-source=csv")
    elif args.data_source == 'sql':
        if not args.db_conn or not args.db_query:
            parser.error("--db-conn and --db-query are required when data-source=sql")
        from src.data_loader import load_sql_query, DataLoaderError
        try:
            df_source = load_sql_query(args.db_conn, args.db_query)
            print(f"✓ Loaded {len(df_source)} rows from SQL query")
        except DataLoaderError as exc:
            parser.error(str(exc))
    elif args.data_source == 'table':
        if not args.db_conn or not args.db_table:
            parser.error("--db-conn and --db-table are required when data-source=table")
        from src.data_loader import load_database_table, DataLoaderError
        try:
            df_source = load_database_table(args.db_conn, args.db_table, args.db_limit)
            print(f"✓ Loaded {len(df_source)} rows from table {args.db_table}")
        except DataLoaderError as exc:
            parser.error(str(exc))
    
    try:
        output_path = generate_report(
            csv_path=args.input,
            output_path=args.out,
            date_start=args.date_start,
            date_end=args.date_end,
            campaigns=args.campaigns,
            use_llm=use_llm,
            use_groq=use_groq,
            use_cache=use_cache,
            dataframe=df_source
        )
        print(f"✓ Report generated: {output_path}")
    except Exception as e:
        logger.error(f"Error generating report: {e}", exc_info=True)
        raise


if __name__ == "__main__":
    main()

