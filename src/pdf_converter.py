"""
PDF conversion utilities for PPTX files.
Uses Python libraries to convert PowerPoint presentations to PDF without requiring LibreOffice.
"""

import os
import logging
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


def convert_pptx_to_pdf_python(
    pptx_path: str,
    pdf_path: Optional[str] = None,
) -> Optional[str]:
    """
    Convert PPTX to PDF using Python libraries (reportlab + Pillow).
    
    This method extracts content from PPTX and creates a PDF directly.
    Falls back to LibreOffice if available, otherwise creates PDF from content.
    
    Args:
        pptx_path: Path to input PPTX file
        pdf_path: Path to output PDF file (default: same as PPTX with .pdf extension)
        
    Returns:
        Path to generated PDF file, or None if conversion failed
    """
    if pdf_path is None:
        pdf_path = pptx_path.replace('.pptx', '.pdf')
    
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.enums import TA_CENTER, TA_LEFT
        from pptx import Presentation
        from PIL import Image as PILImage
        import io
        
        logger.info(f"Converting PPTX to PDF using Python libraries: {pptx_path}")
        
        # Load presentation
        prs = Presentation(pptx_path)
        
        # Create PDF document
        doc = SimpleDocTemplate(pdf_path, pagesize=letter)
        story = []
        styles = getSampleStyleSheet()
        
        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1E88E5'),
            spaceAfter=30,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=18,
            textColor=colors.HexColor('#262730'),
            spaceAfter=12,
            spaceBefore=12
        )
        
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontSize=11,
            textColor=colors.HexColor('#262730'),
            spaceAfter=6,
            leading=14
        )
        
        # Ensure outputs directory exists for temp images
        output_dir = Path(pdf_path).parent
        output_dir.mkdir(parents=True, exist_ok=True)

        temp_files = []

        # Process each slide
        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx > 0:
                story.append(PageBreak())
            
            # Extract text from slide
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            # Add slide content to PDF
            if slide_texts:
                # First text is usually the title
                if len(slide_texts) > 0:
                    first_text = slide_texts[0]
                    if len(first_text) < 100:  # Likely a title
                        story.append(Paragraph(first_text, title_style))
                        story.append(Spacer(1, 0.2*inch))
                        slide_texts = slide_texts[1:]
                
                # Add remaining text
                for text in slide_texts:
                    # Split long text into paragraphs
                    paragraphs = text.split('\n')
                    for para in paragraphs:
                        if para.strip():
                            story.append(Paragraph(para.strip(), normal_style))
                            story.append(Spacer(1, 0.1*inch))
            
            # Try to extract images
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        img = PILImage.open(image_stream)
                        # Resize while maintaining aspect ratio
                        max_width = 6.5 * inch
                        max_height = 4.5 * inch
                        aspect = img.height / img.width if img.width else 1
                        width = max_width
                        height = width * aspect
                        if height > max_height:
                            height = max_height
                            width = height / aspect if aspect else max_width
                        
                        # Save temp image
                        temp_img_path = str(output_dir / f"temp_slide_{slide_idx}_img.png")
                        img.save(temp_img_path)
                        temp_files.append(temp_img_path)
                        
                        # Add to PDF
                        pdf_img = Image(temp_img_path, width=width, height=height)
                        story.append(pdf_img)
                        story.append(Spacer(1, 0.2*inch))
                    except Exception as e:
                        logger.warning(f"Could not extract image from slide {slide_idx}: {e}")
        
        # Build PDF
        doc.build(story)
        
        if Path(pdf_path).exists():
            logger.info(f"PDF generated successfully: {pdf_path}")
            # Clean up temp files
            for temp_file in temp_files:
                try:
                    if os.path.exists(temp_file):
                        os.remove(temp_file)
                except Exception:
                    pass
            return pdf_path
        else:
            logger.error("PDF file was not created")
            return None
            
    except ImportError as e:
        logger.warning(f"Required library not installed: {e}")
        logger.info("Install with: pip install reportlab Pillow")
        return None
    except Exception as e:
        logger.error(f"PDF conversion failed: {e}", exc_info=True)
        return None


def convert_pptx_to_pdf(
    pptx_path: str,
    pdf_path: Optional[str] = None,
) -> Optional[str]:
    """
    Convert PPTX to PDF using best available method.
    
    Tries methods in order:
    1. LibreOffice (if available)
    2. Python libraries (reportlab + Pillow)
    
    Args:
        pptx_path: Path to input PPTX file
        pdf_path: Path to output PDF file
        
    Returns:
        Path to generated PDF file, or None if all methods failed
    """
    if pdf_path is None:
        pdf_path = pptx_path.replace('.pptx', '.pdf')
    
    # Method 1: Try LibreOffice first (fastest, best quality)
    try:
        import subprocess
        result = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf', 
             '--outdir', str(Path(pptx_path).parent), pptx_path],
            capture_output=True,
            timeout=30,
            check=False
        )
        
        if Path(pdf_path).exists():
            logger.info(f"PDF generated using LibreOffice: {pdf_path}")
            return pdf_path
    except (FileNotFoundError, subprocess.TimeoutExpired, Exception) as e:
        logger.info(f"LibreOffice conversion not available: {e}")
    
    # Method 2: Try Python libraries
    logger.info("Attempting PDF conversion using Python libraries...")
    return convert_pptx_to_pdf_python(pptx_path, pdf_path)

